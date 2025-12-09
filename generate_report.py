from pptx import Presentation
from pptx.util import Inches
import pandas as pd
from pathlib import Path

DATA_PATH = Path("data") / "Unified_GBD_Fact_Table_CLEAN.csv"

def generate_basic_ppt(output_path="GBD_Report.pptx"):
    df = pd.read_csv(DATA_PATH)

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Unified GBD Dashboard – National Overview"
    slide.placeholders[1].text = "Automated report generated from Streamlit dataset."

    # Example: add a slide with top 10 causes by DALY in latest year
    latest_year = df["year"].max()
    df_latest = df[df["year"] == latest_year]

    top = (
        df_latest.groupby("cause_name", as_index=False)["val"]
        .sum()
        .sort_values("val", ascending=False)
        .head(10)
    )

    bullet_slide_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    slide2.shapes.title.text = f"Top 10 Causes by DALY Rate – {latest_year}"
    body = slide2.placeholders[1].text_frame

    for _, row in top.iterrows():
        body.add_paragraph().text = f"{row['cause_name']}: {row['val']:.2f}"

    prs.save(output_path)
    print(f"Saved report to {output_path}")

if __name__ == "__main__":
    generate_basic_ppt()
